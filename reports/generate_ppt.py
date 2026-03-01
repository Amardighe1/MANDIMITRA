#!/usr/bin/env python3
"""
MANDIMITRA - Final Year Capstone PPT Generator
================================================
Generates a professional 13-slide PowerPoint for MSBTE Capstone (316004).
Visual-heavy: diagrams, flowcharts, figures on every slide. Minimal text.
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Paths ───────────────────────────────────────────────────
PROJECT_DIR = Path(__file__).parent.parent
REPORTS_DIR = Path(__file__).parent
DIAG        = REPORTS_DIR / "diagrams"
LOGO        = PROJECT_DIR / "download.jpg"
OUTPUT      = REPORTS_DIR / "MANDIMITRA_Capstone_Presentation.pptx"

# ─── Constants ───────────────────────────────────────────────
COLLEGE_LINE1 = "K.E. Society's"
COLLEGE_LINE2 = "Rajarambapu Institute of Technology (Polytechnic)"
COLLEGE_LINE3 = "Lohegaon, Pune - 47"
DEPT_NAME     = "Artificial Intelligence and Machine Learning"
YEAR          = "2025-2026"
TITLE         = "MANDIMITRA"
SUBTITLE      = "Maharashtra Agricultural Market Intelligence\n& Veterinary Services Platform"
COURSE_CODE   = "316004"

STUDENTS = [
    ("Avishkar Borate",      "01"),
    ("Amol Salunke",         "02"),
    ("Pruthviraj Hippirkar", "03"),
    ("Krushna Mane",         "04"),
]
GUIDE     = "Prof. Mayur Gund"
HOD_NAME  = "Prof. Vikram Saste"
PRINCIPAL = "Dr. Kashinath Munde"

# ─── Colors ──────────────────────────────────────────────────
DARK_GREEN   = RGBColor(0x1B, 0x5E, 0x20)
MEDIUM_GREEN = RGBColor(0x2E, 0x7D, 0x32)
LIGHT_GREEN  = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_GREEN = RGBColor(0xA5, 0xD6, 0xA7)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK        = RGBColor(0x00, 0x00, 0x00)
DARK_GREY    = RGBColor(0x33, 0x33, 0x33)
MID_GREY     = RGBColor(0x66, 0x66, 0x66)
LIGHT_GREY   = RGBColor(0xF5, 0xF5, 0xF5)
BLUE         = RGBColor(0x15, 0x65, 0xC0)
ORANGE       = RGBColor(0xE6, 0x51, 0x00)
RED          = RGBColor(0xC6, 0x28, 0x28)
GOLD         = RGBColor(0xFF, 0x8F, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ─── Helper Functions ────────────────────────────────────────
def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, line_color=None):
    """Add a rectangle shape with fill."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK_GREY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", italic=False):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = alignment
    return txBox

def add_multiline_text(slide, left, top, width, height, lines, font_size=16,
                       color=DARK_GREY, bold=False, line_spacing=1.3,
                       alignment=PP_ALIGN.LEFT):
    """Add a text box with multiple paragraphs (bullet-style)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.alignment = alignment
        p.space_after = Pt(font_size * 0.4)
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_bullet_points(slide, left, top, width, height, bullets, font_size=15,
                      color=DARK_GREY, icon_color=MEDIUM_GREEN):
    """Add bullet points with green circle markers."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # Use filled circle as bullet
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0
        # Add bullet marker
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        from lxml import etree
        buNone_list = pPr.findall(qn('a:buNone'))
        for bn in buNone_list:
            pPr.remove(bn)
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '\u25cf')  # filled circle
        buClr = etree.SubElement(pPr, qn('a:buClr'))
        srgb = etree.SubElement(buClr, qn('a:srgbClr'))
        srgb.set('val', f'{icon_color.theme_color if hasattr(icon_color, "theme_color") else "2E7D32"}')

        # Indent
        pPr.set('marL', str(Emu(Inches(0.4))))
        pPr.set('indent', str(Emu(Inches(-0.25))))

    return txBox

def add_image_safe(slide, img_name, left, top, width=None, height=None):
    """Safely add image, returns True if added."""
    path = DIAG / img_name
    if path.exists():
        kwargs = {}
        if width: kwargs['width'] = width
        if height: kwargs['height'] = height
        slide.shapes.add_picture(str(path), left, top, **kwargs)
        return True
    return False

def add_image_caption(slide, text, left, top, width):
    """Add an italicized caption below an image."""
    add_text_box(slide, left, top, width, Inches(0.3), text,
                 font_size=10, color=MID_GREY, italic=True,
                 alignment=PP_ALIGN.CENTER)

def add_slide_number(slide, num, total=13):
    """Add slide number in bottom-right."""
    add_text_box(slide, Inches(12.3), Inches(7.05), Inches(0.9), Inches(0.35),
                 f"{num}/{total}", font_size=10, color=MID_GREY,
                 alignment=PP_ALIGN.RIGHT)

def add_top_bar(slide):
    """Add green top accent bar."""
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), DARK_GREEN)

def add_bottom_bar(slide):
    """Add subtle bottom bar with college name."""
    add_shape(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35), LIGHT_GREY)
    add_text_box(slide, Inches(0.4), Inches(7.15), Inches(8), Inches(0.35),
                 f"{COLLEGE_LINE1} {COLLEGE_LINE2}, {COLLEGE_LINE3}",
                 font_size=8, color=MID_GREY, alignment=PP_ALIGN.LEFT)

def add_section_title(slide, text, left=Inches(0.6), top=Inches(0.25)):
    """Add the slide title with green accent."""
    # Green accent line
    add_shape(slide, left, top + Inches(0.45), Inches(0.06), Inches(0.35), MEDIUM_GREEN)
    add_text_box(slide, left + Inches(0.15), top, Inches(9), Inches(0.8), text,
                 font_size=28, color=DARK_GREEN, bold=True)

def add_stat_card(slide, left, top, width, height, number, label, color=DARK_GREEN):
    """Add a statistics card with big number + label."""
    card = add_rounded_rect(slide, left, top, width, height, WHITE)
    card.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    card.line.width = Pt(1)
    # Number
    add_text_box(slide, left + Inches(0.1), top + Inches(0.1),
                 width - Inches(0.2), Inches(0.6), number,
                 font_size=26, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Label
    add_text_box(slide, left + Inches(0.1), top + Inches(0.6),
                 width - Inches(0.2), Inches(0.4), label,
                 font_size=11, color=MID_GREY, alignment=PP_ALIGN.CENTER)

def add_speaker_notes(slide, text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ══════════════════════════════════════════════════════════════
# SLIDE 1 - TITLE SLIDE
# ══════════════════════════════════════════════════════════════
def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, WHITE)

    # Full-width green header band
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(2.8), DARK_GREEN)

    # Decorative lighter green accent strip
    add_shape(slide, Inches(0), Inches(2.8), SLIDE_W, Inches(0.06), LIGHT_GREEN)

    # Logo on left
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(0.6), Inches(0.35),
                                 width=Inches(1.6))

    # College Name (in white on green band)
    add_text_box(slide, Inches(2.5), Inches(0.25), Inches(10), Inches(0.4),
                 f"{COLLEGE_LINE1} {COLLEGE_LINE2}",
                 font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, Inches(2.5), Inches(0.6), Inches(10), Inches(0.35),
                 COLLEGE_LINE3,
                 font_size=14, color=ACCENT_GREEN, alignment=PP_ALIGN.LEFT)

    # Department and course code
    add_text_box(slide, Inches(2.5), Inches(1.0), Inches(10), Inches(0.35),
                 f"Department of {DEPT_NAME}  |  Course Code: {COURSE_CODE}",
                 font_size=12, color=ACCENT_GREEN, alignment=PP_ALIGN.LEFT)

    # Project Title - Big and centered in green band
    add_text_box(slide, Inches(2.5), Inches(1.45), Inches(10), Inches(0.6),
                 TITLE,
                 font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, Inches(2.5), Inches(2.0), Inches(10), Inches(0.7),
                 SUBTITLE,
                 font_size=16, color=ACCENT_GREEN, alignment=PP_ALIGN.LEFT)

    # Student names section
    add_text_box(slide, Inches(0.6), Inches(3.2), Inches(5), Inches(0.35),
                 "SUBMITTED BY", font_size=11, color=MID_GREY, bold=True)

    y = 3.6
    for name, roll in STUDENTS:
        add_text_box(slide, Inches(0.6), Inches(y), Inches(5), Inches(0.35),
                     f"{name}  (Roll No. {roll})",
                     font_size=14, color=DARK_GREY, bold=False)
        y += 0.35

    # Guide section
    add_text_box(slide, Inches(7.5), Inches(3.2), Inches(5), Inches(0.35),
                 "UNDER THE GUIDANCE OF", font_size=11, color=MID_GREY, bold=True)
    add_text_box(slide, Inches(7.5), Inches(3.6), Inches(5), Inches(0.35),
                 GUIDE, font_size=16, color=MEDIUM_GREEN, bold=True)
    add_text_box(slide, Inches(7.5), Inches(4.0), Inches(5), Inches(0.35),
                 f"HOD: {HOD_NAME}", font_size=12, color=DARK_GREY)
    add_text_box(slide, Inches(7.5), Inches(4.35), Inches(5), Inches(0.35),
                 f"Principal: {PRINCIPAL}", font_size=12, color=DARK_GREY)

    # Academic year
    add_text_box(slide, Inches(7.5), Inches(4.9), Inches(5), Inches(0.35),
                 f"Academic Year: {YEAR}", font_size=12, color=BLUE, bold=True)

    # National thrust area badge
    badge = add_rounded_rect(slide, Inches(0.6), Inches(5.5), Inches(4.5), Inches(0.7), RGBColor(0xE8, 0xF5, 0xE9))
    add_text_box(slide, Inches(0.7), Inches(5.55), Inches(4.3), Inches(0.6),
                 "National Thrust: Digitization | Agriculture | Rural Development",
                 font_size=11, color=MEDIUM_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

    # Thin bottom accent
    add_shape(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35), LIGHT_GREY)
    add_text_box(slide, Inches(0.4), Inches(7.15), Inches(12), Inches(0.35),
                 "Capstone Project Report  |  Diploma in AI & ML  |  MSBTE",
                 font_size=9, color=MID_GREY, alignment=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "Good morning respected jury members, guide, and fellow students. "
        "We are students from the AIML department presenting our capstone project MANDIMITRA. "
        "MANDIMITRA means 'Market Friend' in Marathi. It is a comprehensive agricultural "
        "market intelligence and veterinary services platform built specifically for Maharashtra. "
        "This project falls under the national thrust areas of Digitization, Agriculture modernization, "
        "and Rural Development. Let me introduce our team...")


# ══════════════════════════════════════════════════════════════
# SLIDE 2 - INTRODUCTION
# ══════════════════════════════════════════════════════════════
def slide_introduction(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_section_title(slide, "Introduction")
    add_slide_number(slide, 2)

    # Left side - Key facts as stat cards
    add_text_box(slide, Inches(0.6), Inches(1.1), Inches(6), Inches(0.35),
                 "Maharashtra Agriculture - Key Challenges",
                 font_size=14, color=MEDIUM_GREEN, bold=True)

    # Stat cards row 1
    add_stat_card(slide, Inches(0.6), Inches(1.6), Inches(2.8), Inches(1.0),
                  "1.5 Crore+", "Maharashtra Farmers", DARK_GREEN)
    add_stat_card(slide, Inches(3.6), Inches(1.6), Inches(2.8), Inches(1.0),
                  "15-25%", "Revenue Loss Due to\nPrice Info Asymmetry", RED)

    # Stat cards row 2
    add_stat_card(slide, Inches(0.6), Inches(2.8), Inches(2.8), Inches(1.0),
                  "400+", "APMC Mandis in\nMaharashtra", BLUE)
    add_stat_card(slide, Inches(3.6), Inches(2.8), Inches(2.8), Inches(1.0),
                  "1 : 50,000", "Vet Doctor to\nLivestock Ratio", ORANGE)

    # Problem highlights
    add_text_box(slide, Inches(0.6), Inches(4.1), Inches(6), Inches(0.35),
                 "Two Critical Gaps Identified:",
                 font_size=13, color=DARK_GREEN, bold=True)

    problems = [
        "\u25cf  No Maharashtra-specific ML-based price intelligence platform",
        "\u25cf  No integrated veterinary service with booking & emergency SOS",
        "\u25cf  Existing platforms (e-NAM, MSAMB) lack prediction & advisory",
        "\u25cf  Rural vets are scarce; no verification or discovery system",
    ]
    add_multiline_text(slide, Inches(0.6), Inches(4.5), Inches(6), Inches(2.5),
                       problems, font_size=13, color=DARK_GREY)

    # Right side - DFD Level 0 diagram
    add_text_box(slide, Inches(7.0), Inches(1.1), Inches(6), Inches(0.35),
                 "System Context (DFD Level 0)",
                 font_size=13, color=MEDIUM_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_image_safe(slide, "dfd_level0.png", Inches(7.0), Inches(1.5), width=Inches(5.8))
    add_image_caption(slide, "Fig: Data Flow Diagram - External Entities & Data Stores",
                      Inches(7.0), Inches(6.2), Inches(5.8))

    add_speaker_notes(slide,
        "Maharashtra has over 1.5 crore farmers across 36 districts. Despite having 400+ APMC mandis, "
        "farmers lose 15-25% revenue annually due to price information asymmetry - they don't know "
        "the best time or market to sell. Additionally, rural Maharashtra has a critical shortage of "
        "veterinary doctors - roughly 1 per 50,000 livestock. "
        "No existing platform combines agricultural market intelligence WITH veterinary services. "
        "E-NAM is pan-India with no ML. MSAMB is static data only. "
        "This is the gap MANDIMITRA fills. As the DFD shows, our system interacts with farmers, "
        "doctors, admins, and external data sources like Data.gov.in and NASA POWER.")


# ══════════════════════════════════════════════════════════════
# SLIDE 3 - LITERATURE SURVEY
# ══════════════════════════════════════════════════════════════
def slide_literature(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_section_title(slide, "Literature Survey")
    add_slide_number(slide, 3)

    # Comparison table (as shapes)
    headers = ["System", "ML Price\nForecast", "Crop Risk\nAdvisory", "Vet\nServices", "MH\nSpecific"]
    systems = [
        ("e-NAM",         "No",  "No",  "No",  "No"),
        ("AgriMarket",    "No",  "No",  "No",  "No"),
        ("MSAMB Portal",  "No",  "No",  "No",  "Yes"),
        ("CropIn",        "Ltd", "Ltd", "No",  "No"),
        ("Kisan Suvidha", "No",  "No",  "No",  "No"),
        ("MANDIMITRA",    "Yes", "Yes", "Yes", "Yes"),
    ]

    table_left = Inches(0.5)
    table_top = Inches(1.3)
    col_widths = [Inches(1.8), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.3)]

    # Header row
    x = table_left
    for i, h in enumerate(headers):
        shape = add_shape(slide, x, table_top, col_widths[i], Inches(0.55), DARK_GREEN)
        shape.text_frame.word_wrap = True
        p = shape.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(10)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        shape.text_frame.paragraphs[0].space_before = Pt(2)
        x += col_widths[i] + Inches(0.02)

    # Data rows
    for ri, (sys_name, *vals) in enumerate(systems):
        x = table_left
        row_top = table_top + Inches(0.57) + Inches(ri * 0.42)
        is_mandimitra = sys_name == "MANDIMITRA"
        bg = RGBColor(0xE8, 0xF5, 0xE9) if is_mandimitra else (LIGHT_GREY if ri % 2 == 0 else WHITE)

        for ci, val in enumerate([sys_name] + list(vals)):
            shape = add_shape(slide, x, row_top, col_widths[ci], Inches(0.40), bg)
            p = shape.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(10)
            p.font.bold = is_mandimitra
            p.alignment = PP_ALIGN.CENTER
            if val == "Yes":
                p.font.color.rgb = MEDIUM_GREEN
            elif val == "No":
                p.font.color.rgb = RED
            elif is_mandimitra:
                p.font.color.rgb = DARK_GREEN
            else:
                p.font.color.rgb = DARK_GREY
            x += col_widths[ci] + Inches(0.02)

    # Key research papers (right side)
    add_text_box(slide, Inches(7.5), Inches(1.1), Inches(5.5), Inches(0.35),
                 "Key Research Papers Influencing Design",
                 font_size=13, color=MEDIUM_GREEN, bold=True)

    papers = [
        "1. NeuralCrop (2025) - GDD, VPD, Drought\n   Stress Index features for crop prediction",
        "2. MT-CYP-Net (2025) - Focal-loss inspired\n   class weighting for imbalanced agri data",
        "3. LightGBM (NeurIPS 2017) - Core gradient\n   boosting algorithm for both our models",
        "4. Conformal Prediction (NeurIPS 2019) -\n   Calibrated uncertainty intervals for prices",
        "5. TabPFN (NeurIPS 2022) - Small-sample\n   bootstrapping for minority class training",
    ]
    add_multiline_text(slide, Inches(7.5), Inches(1.5), Inches(5.5), Inches(3.5),
                       papers, font_size=11, color=DARK_GREY, line_spacing=1.2)

    # Comparison chart image
    add_image_safe(slide, "lit_survey_comparison.png", Inches(7.2), Inches(4.1), width=Inches(5.8))
    add_image_caption(slide, "Fig: Comparative Analysis with Existing Systems",
                      Inches(7.2), Inches(6.6), Inches(5.8))

    add_speaker_notes(slide,
        "We conducted a thorough literature survey of 6 existing platforms and 7 key research papers. "
        "As the comparison table shows, NO existing platform provides all four capabilities: "
        "ML price forecasting, crop risk advisory, veterinary services, AND Maharashtra-specific focus. "
        "Our ML design was heavily influenced by recent NeurIPS and ICML papers. "
        "From NeuralCrop 2025, we adopted physics-informed features like GDD and VPD. "
        "From MT-CYP-Net, we used focal-loss inspired class weighting for handling imbalanced risk classes. "
        "The comparative chart shows MANDIMITRA achieves R-squared 0.93 for price prediction "
        "and 87.4% crop risk accuracy - both competitive with or exceeding published results.")


# ══════════════════════════════════════════════════════════════
# SLIDE 4 - PROBLEM STATEMENT & OBJECTIVES
# ══════════════════════════════════════════════════════════════
def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_section_title(slide, "Problem Statement & Objectives")
    add_slide_number(slide, 4)

    # Problem statement box
    prob_box = add_rounded_rect(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.3),
                                 RGBColor(0xFF, 0xEB, 0xEE))
    prob_box.line.color.rgb = RGBColor(0xEF, 0x9A, 0x9A)
    prob_box.line.width = Pt(1.5)

    add_text_box(slide, Inches(0.7), Inches(1.25), Inches(1.5), Inches(0.3),
                 "PROBLEM", font_size=11, color=RED, bold=True)

    prob_text = (
        "Maharashtra's 1.5 crore farmers face 15-25% revenue loss due to price information asymmetry, "
        "and critically limited access to verified veterinary services (1 vet per 50,000 livestock). "
        "No existing platform combines ML-based agricultural market intelligence with "
        "integrated veterinary services for Maharashtra."
    )
    add_text_box(slide, Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.9),
                 prob_text, font_size=13, color=DARK_GREY)

    # 6 Objectives as cards
    add_text_box(slide, Inches(0.5), Inches(2.7), Inches(3), Inches(0.35),
                 "SIX PROJECT OBJECTIVES", font_size=13, color=DARK_GREEN, bold=True)

    objectives = [
        ("OBJ 1", "Data Pipeline", "6.1M+ mandi records, DuckDB\ndedup, Pandera validation,\nMH-only hard constraint"),
        ("OBJ 2", "Weather Integration", "NASA POWER (10yr history) +\nOpen-Meteo (16-day forecast)\nfor all 36 MH districts"),
        ("OBJ 3", "Crop Risk Model", "LightGBM + GDD/VPD/DSI\nphysics features, 87.4% acc,\nfocal-loss class weighting"),
        ("OBJ 4", "Price Forecasting", "5-horizon LightGBM, R2=0.93,\nconformal prediction intervals\nat 80/90/95% confidence"),
        ("OBJ 5", "Web Application", "Next.js 14 + FastAPI + Supabase\n3 role-based dashboards,\nmodern UI with animations"),
        ("OBJ 6", "Vet Services", "15 API endpoints, doctor\nverification, booking, and\nemergency SOS (FCFS)"),
    ]

    for i, (num, title, desc) in enumerate(objectives):
        col = i % 3
        row = i // 3
        left = Inches(0.5) + col * Inches(4.15)
        top = Inches(3.1) + row * Inches(2.0)

        card = add_rounded_rect(slide, left, top, Inches(3.95), Inches(1.85),
                                 RGBColor(0xE8, 0xF5, 0xE9))
        # Number badge
        badge = add_rounded_rect(slide, left + Inches(0.1), top + Inches(0.1),
                                  Inches(0.7), Inches(0.35), MEDIUM_GREEN)
        badge_tf = badge.text_frame
        badge_tf.paragraphs[0].text = num
        badge_tf.paragraphs[0].font.size = Pt(9)
        badge_tf.paragraphs[0].font.color.rgb = WHITE
        badge_tf.paragraphs[0].font.bold = True
        badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Title
        add_text_box(slide, left + Inches(0.9), top + Inches(0.08),
                     Inches(2.9), Inches(0.35), title,
                     font_size=13, color=DARK_GREEN, bold=True)
        # Description
        add_text_box(slide, left + Inches(0.15), top + Inches(0.5),
                     Inches(3.65), Inches(1.3), desc,
                     font_size=11, color=DARK_GREY)

    add_speaker_notes(slide,
        "The problem is clear: farmers lose money because they don't have access to reliable "
        "price intelligence, and they lose livestock because they can't find verified vets quickly. "
        "We defined 6 specific, measurable objectives. Objective 1 is the data foundation - "
        "processing 6.1 million mandi records. Objectives 2-4 are the ML intelligence layer. "
        "Objectives 5-6 are the user-facing application. Each objective has clear success metrics "
        "that we will show in the Results slide.")


# ══════════════════════════════════════════════════════════════
# SLIDE 5 - SCOPE OF THE PROJECT
# ══════════════════════════════════════════════════════════════
def slide_scope(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_section_title(slide, "Scope of the Project")
    add_slide_number(slide, 5)

    # Scope diagram (the main visual)
    add_image_safe(slide, "scope_diagram.png", Inches(0.4), Inches(1.2), width=Inches(6.5))
    add_image_caption(slide, "Fig: In-Scope vs Out-of-Scope Features",
                      Inches(0.4), Inches(6.5), Inches(6.5))

    # Right side - scope details
    add_text_box(slide, Inches(7.2), Inches(1.2), Inches(5.5), Inches(0.35),
                 "IN SCOPE", font_size=14, color=MEDIUM_GREEN, bold=True)

    in_scope = [
        "\u2705  Maharashtra state (36 districts, 400+ mandis)",
        "\u2705  3 User Roles: Farmer, Doctor, Admin",
        "\u2705  ML: Crop Risk Advisory + Price Forecasting",
        "\u2705  Vet: Verification, Booking, Emergency SOS",
        "\u2705  6.1M mandi records + 10 years weather data",
        "\u2705  Full-stack web app (Next.js + FastAPI + Supabase)",
    ]
    add_multiline_text(slide, Inches(7.2), Inches(1.6), Inches(5.8), Inches(2.5),
                       in_scope, font_size=12, color=DARK_GREY)

    add_text_box(slide, Inches(7.2), Inches(4.2), Inches(5.5), Inches(0.35),
                 "OUT OF SCOPE (Current Version)", font_size=14, color=RED, bold=True)

    out_scope = [
        "\u274c  Other Indian states (MH-only hard constraint)",
        "\u274c  Direct e-commerce / buy-sell of produce",
        "\u274c  Chatbot / conversational AI interface",
        "\u274c  Native mobile app (React Native planned)",
        "\u274c  Marathi language support (future Phase 2)",
        "\u274c  Pesticide / fertilizer recommendations",
    ]
    add_multiline_text(slide, Inches(7.2), Inches(4.6), Inches(5.8), Inches(2.5),
                       out_scope, font_size=12, color=DARK_GREY)

    add_speaker_notes(slide,
        "The scope of MANDIMITRA is strictly limited to Maharashtra state. We enforce this "
        "as a hard constraint in the data pipeline using the special 'filters[state.keyword]' "
        "parameter in the Data.gov.in API. This ensures zero contamination from other states. "
        "We cover all 36 districts with 400+ mandis. We support 3 user roles with distinct dashboards. "
        "The current version is a web-only application. We've identified a clear Phase 2 roadmap "
        "including React Native mobile app and Marathi language support for wider adoption.")


# ══════════════════════════════════════════════════════════════
# SLIDE 6 - FEASIBILITY ANALYSIS
# ══════════════════════════════════════════════════════════════
def slide_feasibility(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_section_title(slide, "Feasibility Analysis")
    add_slide_number(slide, 6)

    # Four feasibility cards in 2x2 grid
    cards = [
        ("Technical Feasibility", MEDIUM_GREEN,
         ["Python 3.10+ / Node.js 18+ (mature, well-supported)",
          "LightGBM: NeurIPS-published, proven on tabular data",
          "AGMARKNET open API + Kaggle: data freely available",
          "Supabase: managed PostgreSQL + Auth + Storage"]),
        ("Financial Feasibility", BLUE,
         ["Development Cost: Rs 0 (all open-source tools)",
          "Production Cost: ~Rs 6,800/yr (VPS + domain)",
          "Supabase Free Tier: 500MB DB, 1GB storage",
          "Vercel Free Tier: Next.js hosting included"]),
        ("Operational Feasibility", ORANGE,
         ["Target users: smartphone-owning MH farmers",
          "Simple web UI, no training needed for farmers",
          "Admin verification ensures doctor quality",
          "SOS is single-tap emergency broadcast"]),
        ("Schedule Feasibility", RGBColor(0x6A, 0x1B, 0x9A),
         ["Phase 1: Data Pipeline (4 weeks) - DONE",
          "Phase 2: ML Development (3 weeks) - DONE",
          "Phase 3: Web App (4 weeks) - DONE",
          "Phase 4: Testing & Docs (2 weeks) - DONE"]),
    ]

    for i, (title, color, points) in enumerate(cards):
        col = i % 2
        row = i // 2
        left = Inches(0.5) + col * Inches(6.35)
        top = Inches(1.2) + row * Inches(3.0)

        # Card background
        card = add_rounded_rect(slide, left, top, Inches(6.1), Inches(2.8), WHITE)
        card.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        card.line.width = Pt(1)

        # Color top strip
        add_shape(slide, left, top, Inches(6.1), Inches(0.06), color)

        # Title
        add_text_box(slide, left + Inches(0.2), top + Inches(0.15),
                     Inches(5.7), Inches(0.4), title,
                     font_size=15, color=color, bold=True)

        # Points
        y_pt = top + Inches(0.6)
        for pt in points:
            add_text_box(slide, left + Inches(0.3), y_pt,
                         Inches(5.6), Inches(0.4), f"\u25cf  {pt}",
                         font_size=11, color=DARK_GREY)
            y_pt += Inches(0.42)

    add_speaker_notes(slide,
        "We analyzed feasibility across 4 dimensions. Technically, we're using proven, "
        "mature frameworks - LightGBM is a NeurIPS-published algorithm, and all data sources are "
        "freely available through government open data. Financially, the entire platform can run on "
        "free-tier services during development, and only ~Rs 6,800/year for production. "
        "Operationally, our target users already have smartphones, and the UI is designed to be "
        "intuitive. Schedule-wise, we completed all 4 phases within 13 weeks as planned.")


# ══════════════════════════════════════════════════════════════
# SLIDE 7 - METHODOLOGY / APPROACH
# ══════════════════════════════════════════════════════════════
def slide_methodology(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_section_title(slide, "Methodology / Approach")
    add_slide_number(slide, 7)

    # Pipeline flowchart - main visual
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(0.35),
                 "End-to-End Data & ML Pipeline (6-Stage Flow)",
                 font_size=14, color=MEDIUM_GREEN, bold=True)
    add_image_safe(slide, "pipeline_flowchart.png", Inches(0.3), Inches(1.5), width=Inches(7.0))
    add_image_caption(slide, "Fig: Complete Pipeline from Raw Data to Deployed App",
                      Inches(0.3), Inches(6.5), Inches(7.0))

    # Right side - 6 stages summary
    add_text_box(slide, Inches(7.6), Inches(1.1), Inches(5.5), Inches(0.35),
                 "Agile Methodology (4 Phases, 13 Weeks)",
                 font_size=13, color=MEDIUM_GREEN, bold=True)

    stages = [
        ("Stage 1", "DATA INGESTION", "4 sources: AGMARKNET, Kaggle,\nNASA POWER, Open-Meteo"),
        ("Stage 2", "VALIDATION", "Pandera schemas, MH-only filter,\ntype & range enforcement"),
        ("Stage 3", "DEDUPLICATION", "7-column natural key, DuckDB,\n~700K duplicates removed"),
        ("Stage 4", "FEATURE ENG.", "GDD, VPD, Drought Index,\nrolling stats, lag features"),
        ("Stage 5", "ML TRAINING", "LightGBM 5-fold CV,\nclass weights, conformal intervals"),
        ("Stage 6", "SERVING", "FastAPI endpoints, .joblib models,\nNext.js frontend display"),
    ]

    for i, (num, title, desc) in enumerate(stages):
        y = Inches(1.55) + i * Inches(0.82)
        # Number circle
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.6), y,
                                       Inches(0.35), Inches(0.35))
        circ.fill.solid()
        circ.fill.fore_color.rgb = MEDIUM_GREEN
        circ.line.fill.background()
        circ.text_frame.paragraphs[0].text = num.split()[1]
        circ.text_frame.paragraphs[0].font.size = Pt(8)
        circ.text_frame.paragraphs[0].font.color.rgb = WHITE
        circ.text_frame.paragraphs[0].font.bold = True
        circ.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_text_box(slide, Inches(8.05), y - Inches(0.02),
                     Inches(1.5), Inches(0.3), title,
                     font_size=11, color=DARK_GREEN, bold=True)
        add_text_box(slide, Inches(9.5), y - Inches(0.02),
                     Inches(3.5), Inches(0.7), desc,
                     font_size=10, color=DARK_GREY)

    add_speaker_notes(slide,
        "We followed Agile methodology with 4 phases over 13 weeks. The pipeline has 6 stages. "
        "Stage 1 ingests data from 4 sources. The critical detail is using 'filters[state.keyword]' "
        "for exact Maharashtra matching in the AGMARKNET API. Stage 2 validates with Pandera schemas. "
        "Stage 3 removes 700K duplicates using a 7-column natural key in DuckDB. "
        "Stage 4 is where we add physics-informed features from the NeuralCrop paper. "
        "Stage 5 trains our LightGBM models with 5-fold cross-validation. "
        "Stage 6 serves predictions through FastAPI REST endpoints consumed by the Next.js frontend.")


# ══════════════════════════════════════════════════════════════
# SLIDE 8 - SYSTEM ARCHITECTURE & DESIGN
# ══════════════════════════════════════════════════════════════
def slide_design1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_section_title(slide, "System Architecture & Design")
    add_slide_number(slide, 8)

    # System architecture - main diagram
    add_image_safe(slide, "system_architecture.png", Inches(0.3), Inches(1.1), width=Inches(6.3))
    add_image_caption(slide, "Fig: Three-Tier System Architecture",
                      Inches(0.3), Inches(5.6), Inches(6.3))

    # Use case diagram - right side
    add_image_safe(slide, "use_cases.png", Inches(6.8), Inches(1.1), width=Inches(6.2))
    add_image_caption(slide, "Fig: Use Case Diagram (12 Use Cases, 3 Actors)",
                      Inches(6.8), Inches(5.6), Inches(6.2))

    # Architecture labels at bottom
    layers = [
        ("PRESENTATION", "Next.js 14 / Tailwind / Framer Motion", BLUE),
        ("APPLICATION", "FastAPI / LightGBM / Auth Middleware", MEDIUM_GREEN),
        ("DATA", "Supabase PostgreSQL / DuckDB / Parquet", ORANGE),
    ]
    for i, (layer, tech, color) in enumerate(layers):
        left = Inches(0.5) + i * Inches(4.3)
        badge = add_rounded_rect(slide, left, Inches(6.0), Inches(4.0), Inches(0.9), WHITE)
        badge.line.color.rgb = color
        badge.line.width = Pt(1.5)
        add_text_box(slide, left + Inches(0.1), Inches(6.0),
                     Inches(3.8), Inches(0.35), layer,
                     font_size=11, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.1), Inches(6.35),
                     Inches(3.8), Inches(0.45), tech,
                     font_size=9, color=DARK_GREY, alignment=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "MANDIMITRA uses a three-tier architecture. The Presentation Layer is Next.js 14 with "
        "server-side rendering, Tailwind CSS, and Framer Motion animations. The Application Layer "
        "has FastAPI with 2 main modules: auth.py and vet.py with 15 endpoints. ML models load "
        "from .joblib files. The Data Layer uses a dual-database strategy: Supabase PostgreSQL "
        "for transactional data and DuckDB for OLAP analytics over 6.1M records. "
        "The use case diagram shows 12 use cases across 3 actors: Farmers have 6, including "
        "SOS broadcast. Doctors have 4, including emergency acceptance. Admins have 2 for oversight.")


# ══════════════════════════════════════════════════════════════
# SLIDE 9 - WORKING PROTOTYPE DETAILS
# ══════════════════════════════════════════════════════════════
def slide_design2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_section_title(slide, "Working Prototype & Key Flows")
    add_slide_number(slide, 9)

    # ML Pipeline diagram - top left
    add_image_safe(slide, "ml_pipeline.png", Inches(0.3), Inches(1.1), width=Inches(6.3))
    add_image_caption(slide, "Fig: ML Training Pipeline (Feature Eng. -> LightGBM -> Conformal)",
                      Inches(0.3), Inches(4.6), Inches(6.3))

    # Vet service flow - top right
    add_image_safe(slide, "vet_service_flow.png", Inches(6.8), Inches(1.1), width=Inches(6.2))
    add_image_caption(slide, "Fig: Veterinary Service Workflow (Swimlane Diagram)",
                      Inches(6.8), Inches(4.6), Inches(6.2))

    # Bottom - ER Diagram + Booking Sequence
    add_image_safe(slide, "er_diagram.png", Inches(0.3), Inches(5.0), width=Inches(4.2))
    add_image_caption(slide, "Fig: ER Diagram (3 Tables + 2 Data Stores)",
                      Inches(0.3), Inches(6.85), Inches(4.2))

    add_image_safe(slide, "booking_sequence.png", Inches(4.7), Inches(5.0), width=Inches(4.2))
    add_image_caption(slide, "Fig: Booking Sequence Diagram",
                      Inches(4.7), Inches(6.85), Inches(4.2))

    add_image_safe(slide, "dedup_process_flow.png", Inches(9.1), Inches(5.0), width=Inches(4.0))
    add_image_caption(slide, "Fig: Deduplication Strategy (7-Key)",
                      Inches(9.1), Inches(6.85), Inches(4.0))

    add_speaker_notes(slide,
        "This slide shows the core working processes of our prototype. "
        "TOP LEFT: The ML pipeline shows our process from raw data through feature engineering "
        "(GDD, VPD, DSI), LightGBM training with 5-fold CV, to conformal calibration for uncertainty. "
        "TOP RIGHT: The veterinary service swimlane shows the complete flow across farmer, "
        "system, doctor, and admin lanes - from registration through verification, booking, "
        "and emergency SOS handling with first-come-first-serve acceptance. "
        "BOTTOM: The ER diagram shows 3 PostgreSQL tables (profiles, bookings, emergency_requests). "
        "The booking sequence shows the full API interaction. And the deduplication flow shows how "
        "we reduce 6.8M raw records to 6.1M using a 7-column natural key in DuckDB.")


# ══════════════════════════════════════════════════════════════
# SLIDE 10 - RESULTS & APPLICATIONS
# ══════════════════════════════════════════════════════════════
def slide_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_section_title(slide, "Results & Applications")
    add_slide_number(slide, 10)

    # ML Results diagrams
    add_image_safe(slide, "crop_risk_results.png", Inches(0.3), Inches(1.1), width=Inches(6.3))
    add_image_caption(slide, "Fig: Crop Risk Advisor - 87.4% Accuracy, Feature Importance",
                      Inches(0.3), Inches(4.35), Inches(6.3))

    add_image_safe(slide, "price_model_results.png", Inches(6.8), Inches(1.1), width=Inches(6.2))
    add_image_caption(slide, "Fig: Price Intelligence - R2=0.93 (1-day), 5 Horizons",
                      Inches(6.8), Inches(4.35), Inches(6.2))

    # Key metrics bar at bottom
    metrics = [
        ("87.4%", "Crop Risk\nAccuracy", MEDIUM_GREEN),
        ("R\u00b2 = 0.93", "1-Day Price\nForecast", BLUE),
        ("6.1M+", "Mandi Records\nProcessed", DARK_GREEN),
        ("35/36", "Districts\nCovered", ORANGE),
        ("15", "Vet API\nEndpoints", RGBColor(0x6A, 0x1B, 0x9A)),
        ("< 500ms", "API Response\nTime", RGBColor(0x00, 0x88, 0x88)),
    ]

    for i, (val, label, color) in enumerate(metrics):
        left = Inches(0.4) + i * Inches(2.1)
        card = add_rounded_rect(slide, left, Inches(4.7), Inches(1.95), Inches(1.15), WHITE)
        card.line.color.rgb = color
        card.line.width = Pt(2)
        add_text_box(slide, left + Inches(0.05), Inches(4.75),
                     Inches(1.85), Inches(0.5), val,
                     font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.05), Inches(5.2),
                     Inches(1.85), Inches(0.5), label,
                     font_size=9, color=MID_GREY, alignment=PP_ALIGN.CENTER)

    # Data coverage + applications
    add_image_safe(slide, "data_coverage.png", Inches(0.3), Inches(6.0), width=Inches(5.5))

    add_text_box(slide, Inches(6.0), Inches(6.0), Inches(7), Inches(0.3),
                 "PRACTICAL APPLICATIONS:", font_size=11, color=DARK_GREEN, bold=True)
    apps = [
        "\u25cf  Farmers: Optimal sell timing, market selection, risk preparedness",
        "\u25cf  Veterinary: Quick verified doctor discovery, emergency response",
        "\u25cf  Government: Data-driven agriculture policy insights",
    ]
    add_multiline_text(slide, Inches(6.0), Inches(6.3), Inches(7), Inches(1.0),
                       apps, font_size=10, color=DARK_GREY)

    add_speaker_notes(slide,
        "Here are our results. The Crop Risk Advisor achieves 87.4% overall accuracy. "
        "The key achievement is improving High Risk recall from 34.9% to 45.9% using "
        "focal-loss inspired weights (1:10:50). The GDD physics feature ranked 9th in importance, "
        "validating the NeuralCrop approach. "
        "The Price Intelligence Engine achieves R-squared 0.93 for 1-day forecasts. Even at "
        "14-day horizon, R-squared stays above 0.86. Conformal intervals give farmers "
        "calibrated uncertainty at 80/90/95% confidence. "
        "We process 6.1M+ records covering 35 of 36 districts. All 15 vet API endpoints pass "
        "testing. API response times are under 500ms. "
        "Practical applications include: farmers can time their sales better, instantly find "
        "verified vets, and send one-tap SOS in emergencies.")


# ══════════════════════════════════════════════════════════════
# SLIDE 11 - CONCLUSION
# ══════════════════════════════════════════════════════════════
def slide_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_section_title(slide, "Conclusion & Future Scope")
    add_slide_number(slide, 11)

    # Left: Achievements
    add_text_box(slide, Inches(0.6), Inches(1.2), Inches(6), Inches(0.35),
                 "KEY ACHIEVEMENTS", font_size=14, color=DARK_GREEN, bold=True)

    achievements = [
        "\u2705  Built production-quality pipeline: 6.1M records, 4 sources",
        "\u2705  Crop Risk Advisor: 87.4% acc with physics-informed features",
        "\u2705  Price Intelligence: R\u00b2=0.93 with conformal prediction intervals",
        "\u2705  Full-stack web app: 3 dashboards, modern responsive UI",
        "\u2705  Veterinary module: 15 endpoints, booking + emergency SOS",
        "\u2705  All 6 objectives met with measurable success metrics",
    ]
    add_multiline_text(slide, Inches(0.6), Inches(1.6), Inches(6), Inches(2.8),
                       achievements, font_size=13, color=DARK_GREY)

    # Technical contributions
    add_text_box(slide, Inches(0.6), Inches(4.2), Inches(6), Inches(0.35),
                 "TECHNICAL CONTRIBUTIONS", font_size=14, color=BLUE, bold=True)

    contributions = [
        "\u25cf  First: Physics-informed ML (GDD/VPD) for MH crop risk",
        "\u25cf  First: Conformal prediction for agricultural price forecasts",
        "\u25cf  First: Unified agri-intelligence + vet platform for MH",
        "\u25cf  Open-source stack: Rs 0 development cost",
    ]
    add_multiline_text(slide, Inches(0.6), Inches(4.6), Inches(6), Inches(2.0),
                       contributions, font_size=12, color=DARK_GREY)

    # Right: Future scope + deployment
    add_text_box(slide, Inches(7.0), Inches(1.2), Inches(6), Inches(0.35),
                 "FUTURE SCOPE (Phase 2-3)", font_size=14, color=ORANGE, bold=True)

    future = [
        "\u25b6  React Native mobile app (Android + iOS)",
        "\u25b6  Marathi language support for rural adoption",
        "\u25b6  iTransformer for improved long-horizon forecasts",
        "\u25b6  Real-time price alerts via WebSocket push",
        "\u25b6  GPS-based nearest vet doctor discovery",
        "\u25b6  Telemedicine video consultations",
        "\u25b6  Crop disease identification via phone camera",
    ]
    add_multiline_text(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(2.8),
                       future, font_size=12, color=DARK_GREY)

    # Deployment architecture
    add_image_safe(slide, "deployment_architecture.png", Inches(7.0), Inches(4.2), width=Inches(5.8))
    add_image_caption(slide, "Fig: Production Deployment Architecture",
                      Inches(7.0), Inches(6.6), Inches(5.8))

    add_speaker_notes(slide,
        "In conclusion, MANDIMITRA successfully achieves all 6 objectives. We built a "
        "complete end-to-end system from data pipeline to web application. "
        "Our 3 key technical contributions are: applying physics-informed ML features to "
        "Maharashtra agriculture for the first time, using conformal prediction for farmer-facing "
        "price uncertainty, and creating the first unified agricultural + veterinary platform for MH. "
        "For future work, we plan to build a React Native mobile app and add Marathi language "
        "support, which are critical for actual deployment in rural Maharashtra. "
        "We also plan to explore transformer-based models like iTransformer for improved "
        "long-horizon price forecasting accuracy. The deployment architecture shows how the "
        "current system can be deployed using Vercel and a VPS for under Rs 7000 per year.")


# ══════════════════════════════════════════════════════════════
# SLIDE 12 - REFERENCES
# ══════════════════════════════════════════════════════════════
def slide_references(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_section_title(slide, "References")
    add_slide_number(slide, 12)

    refs_left = [
        "[1]  Data.gov.in - AGMARKNET Agricultural Market\n      Prices, Govt. of India, 2001-2026",
        "[2]  NASA POWER - Worldwide Energy Resources,\n      NASA Langley Research Center",
        "[3]  Open-Meteo - Free Weather Forecast API",
        "[4]  Kaggle - India Agri Commodity Prices Archive",
        "[5]  G. Ke et al., \"LightGBM: Highly Efficient GBDT\",\n      NeurIPS 2017",
        "[6]  Patel et al., \"MT-CYP-Net: Multi-Task Crop\n      Yield Prediction\", arXiv:2505.12069, 2025",
        "[7]  Kumar et al., \"NeuralCrop: Physics + ML for\n      Crop Yield\", arXiv:2512.20177, 2025",
        "[8]  Romano et al., \"Conformalized Quantile\n      Regression\", NeurIPS 2019",
        "[9]  Hollmann et al., \"TabPFN: Transformer for\n      Tabular Data\", NeurIPS 2022",
        "[10] Weber et al., \"Sub-Field Crop Yield Prediction\n       Explainability\", ICML 2024",
    ]

    refs_right = [
        "[11] Supabase - Open Source Firebase Alternative\n       https://supabase.com/",
        "[12] Next.js 14 - React Framework for Production\n       https://nextjs.org/",
        "[13] FastAPI - Modern Python Web Framework\n       https://fastapi.tiangolo.com/",
        "[14] Tailwind CSS - Utility-First CSS Framework\n       https://tailwindcss.com/",
        "[15] DuckDB - In-Process OLAP Database\n       https://duckdb.org/",
        "[16] Framer Motion - React Animation Library\n       https://www.framer.com/motion/",
        "[17] MSAMB - MH State Agricultural Marketing\n       Board, https://www.msamb.com/",
        "[18] e-NAM - National Agriculture Market\n       https://www.enam.gov.in/",
        "[19] 20th Livestock Census, GoI, 2019",
        "[20] Sharma et al., \"Intrinsic Explainability\",\n       arXiv:2508.06939, 2025",
    ]

    add_multiline_text(slide, Inches(0.5), Inches(1.2), Inches(6.2), Inches(6.0),
                       refs_left, font_size=10, color=DARK_GREY, line_spacing=1.15)
    add_multiline_text(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(6.0),
                       refs_right, font_size=10, color=DARK_GREY, line_spacing=1.15)

    add_speaker_notes(slide,
        "These are our 20 references. Key highlights: References 5-10 are the NeurIPS/ICML/arXiv "
        "research papers that directly influenced our ML model design. References 1-4 are our data "
        "sources. References 11-16 are the core technology frameworks we used. "
        "All references are cited in the full project report.")


# ══════════════════════════════════════════════════════════════
# SLIDE 13 - THANK YOU / Q&A
# ══════════════════════════════════════════════════════════════
def slide_thankyou(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Green background
    set_slide_bg(slide, DARK_GREEN)

    # Decorative lighter strip
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), LIGHT_GREEN)
    add_shape(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), LIGHT_GREEN)

    # Main text
    centered_text = slide.shapes.add_textbox(
        Inches(1), Inches(1.5), Inches(11.333), Inches(1.0))
    tf = centered_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(48)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    add_text_box(slide, Inches(1), Inches(2.6), Inches(11.333), Inches(0.5),
                 "Questions & Discussion", font_size=24, color=ACCENT_GREEN,
                 alignment=PP_ALIGN.CENTER)

    # Divider
    add_shape(slide, Inches(4.5), Inches(3.3), Inches(4.333), Inches(0.03), LIGHT_GREEN)

    # Project details
    add_text_box(slide, Inches(1), Inches(3.6), Inches(11.333), Inches(0.5),
                 TITLE, font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4.1), Inches(11.333), Inches(0.5),
                 "Maharashtra Agricultural Market Intelligence & Veterinary Services Platform",
                 font_size=13, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

    # Team names
    team_text = " | ".join([f"{n} ({r})" for n, r in STUDENTS])
    add_text_box(slide, Inches(1), Inches(5.0), Inches(11.333), Inches(0.4),
                 team_text, font_size=12, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

    # Guide
    add_text_box(slide, Inches(1), Inches(5.5), Inches(11.333), Inches(0.4),
                 f"Guide: {GUIDE}  |  {COLLEGE_LINE1} {COLLEGE_LINE2}",
                 font_size=11, color=RGBColor(0x81, 0xC7, 0x84), alignment=PP_ALIGN.CENTER)

    # Contact/GitHub badge (placeholder)
    add_text_box(slide, Inches(1), Inches(6.3), Inches(11.333), Inches(0.4),
                 f"Department of {DEPT_NAME}  |  Academic Year {YEAR}",
                 font_size=10, color=RGBColor(0x66, 0xBB, 0x6A), alignment=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "Thank you for your attention. We're happy to answer any questions about "
        "the data pipeline, ML models, the web application, or the veterinary service module. "
        "We can also demonstrate the live prototype if needed. Thank you.")


# ══════════════════════════════════════════════════════════════
# BUILD
# ══════════════════════════════════════════════════════════════
def main():
    print("Building MANDIMITRA Capstone Presentation (.pptx)...")
    print(f"  Output: {OUTPUT}")

    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 13.333 inches (widescreen 16:9)
    prs.slide_height = Emu(6858000)   # 7.5 inches

    slide_title(prs)        # Slide 1
    slide_introduction(prs) # Slide 2
    slide_literature(prs)   # Slide 3
    slide_problem(prs)      # Slide 4
    slide_scope(prs)        # Slide 5
    slide_feasibility(prs)  # Slide 6
    slide_methodology(prs)  # Slide 7
    slide_design1(prs)      # Slide 8
    slide_design2(prs)      # Slide 9
    slide_results(prs)      # Slide 10
    slide_conclusion(prs)   # Slide 11
    slide_references(prs)   # Slide 12
    slide_thankyou(prs)     # Slide 13

    prs.save(str(OUTPUT))
    size = os.path.getsize(OUTPUT)
    print(f"\n  Slides: 13")
    print(f"  Size: {size:,} bytes ({size/1024/1024:.1f} MB)")
    print(f"  Diagrams embedded: system_architecture, dfd_level0, use_cases,")
    print(f"    scope_diagram, pipeline_flowchart, ml_pipeline, vet_service_flow,")
    print(f"    er_diagram, booking_sequence, dedup_process_flow, crop_risk_results,")
    print(f"    price_model_results, data_coverage, lit_survey_comparison,")
    print(f"    deployment_architecture + college logo")
    print(f"\n  Done! Open in PowerPoint or Google Slides.")


if __name__ == "__main__":
    main()
